"""
RiskWaiver360 — PPTX generator
Red + Black Cybersecurity GRC theme
12 slides, speaker notes, professional layout
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy, pathlib, os

OUT = pathlib.Path(__file__).parent / "RiskWaiver360_Final_Presentation.pptx"

# ── palette ──────────────────────────────────────────────────────────────────
BG        = RGBColor(0x02, 0x06, 0x17)   # slide background
CARD      = RGBColor(0x0D, 0x19, 0x29)   # card fill
CARD2     = RGBColor(0x0F, 0x17, 0x2A)   # card fill alt
RED       = RGBColor(0xDC, 0x26, 0x26)   # primary red
RED_LT    = RGBColor(0xEF, 0x44, 0x44)   # lighter red
RED_DIM   = RGBColor(0x99, 0x1B, 0x1B)   # dark red
WHITE     = RGBColor(0xF8, 0xFA, 0xFC)   # body text
MUTED     = RGBColor(0x94, 0xA3, 0xB8)   # muted text
MUTED2    = RGBColor(0x64, 0x74, 0x8B)   # dimmer muted
BORDER    = RGBColor(0x1E, 0x29, 0x3B)   # card border
GREEN     = RGBColor(0x16, 0xA3, 0x4A)
GREEN_LT  = RGBColor(0x86, 0xEF, 0xAC)
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_LT  = RGBColor(0xFC, 0xD3, 0x4D)
PURPLE    = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LT = RGBColor(0xC4, 0xB5, 0xFD)
CYAN      = RGBColor(0x06, 0xB6, 0xD4)
CYAN_LT   = RGBColor(0x67, 0xE8, 0xF9)
RED_ROW   = RGBColor(0x7F, 0x1D, 0x1D)   # table header row
STRIPE    = RGBColor(0x0A, 0x12, 0x20)   # alternate row stripe

# ── slide dimensions (widescreen 16:9) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ═══════════════════════════════════════════════════════════════════════════
# helpers
# ═══════════════════════════════════════════════════════════════════════════

def add_slide():
    sl = prs.slides.add_slide(BLANK)
    fill_bg(sl)
    return sl


def fill_bg(sl):
    bg = sl.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def rgb_hex(r: RGBColor) -> str:
    return f"{r[0]:02X}{r[1]:02X}{r[2]:02X}"


def add_notes(sl, text: str):
    notes = sl.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


def box(sl, x, y, w, h,
        fill=None, border=None, border_w=Pt(1), radius=None):
    """Add a plain rectangle."""
    shape = sl.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.width = border_w if border else Pt(0)
    if border:
        shape.line.color.rgb = border
    else:
        shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.shadow.inherit = False
    return shape


def accent_bar(sl, x, y, w=0.04, h=0.55, color=None):
    """Thin vertical left-accent bar for cards."""
    color = color or RED
    b = box(sl, x, y, w, h, fill=color)
    b.line.fill.background()
    return b


def txt(sl, text, x, y, w, h,
        size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
        italic=False, wrap=True, font="Calibri"):
    color = color or WHITE
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = bold and "Calibri" or font
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return txb


def txt_block(sl, lines, x, y, w, h,
              size=12, color=None, bold_first=False,
              spacing=1.15, font="Calibri", align=PP_ALIGN.LEFT):
    """Multi-line textbox; each line is a separate paragraph."""
    color = color or WHITE
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
    return txb


def label_value(sl, label, value, x, y, lw=1.5, vw=4.5, size=11,
                label_color=None, value_color=None):
    label_color = label_color or MUTED
    value_color = value_color or WHITE
    txt(sl, label, x, y, lw, 0.3, size=size, color=label_color, bold=True)
    txt(sl, value, x + lw, y, vw, 0.3, size=size, color=value_color)


def divider(sl, x, y, w, color=None, thick=Pt(0.75)):
    color = color or RED
    ln = sl.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = thick


def card_box(sl, x, y, w, h, accent=True, accent_color=None, fill=None):
    """Dark glass card with optional red top-accent line."""
    fill = fill or CARD
    accent_color = accent_color or RED
    c = box(sl, x, y, w, h, fill=fill, border=BORDER, border_w=Pt(0.75))
    if accent:
        # 3-pt red top accent line
        top = box(sl, x, y, w, 0.045, fill=accent_color)
        top.line.fill.background()
    return c


def section_header(sl, number, title, subtitle=None):
    """Standard slide header: section number badge + title + subtitle."""
    # section number pill
    badge = box(sl, 0.35, 0.25, 0.45, 0.42, fill=RED)
    badge.line.fill.background()
    txt(sl, str(number), 0.35, 0.26, 0.45, 0.38,
        size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(sl, title, 0.9, 0.25, 11.5, 0.5,
        size=26, bold=True, color=WHITE)

    if subtitle:
        txt(sl, subtitle, 0.9, 0.72, 11.5, 0.3,
            size=13, color=MUTED, italic=True)

    divider(sl, 0.35, 0.95, 12.65, color=RED, thick=Pt(1.2))


def table_shape(sl, headers, rows,
                x, y, w, h,
                col_widths=None,
                hdr_color=None,
                hdr_text_color=None,
                stripe_color=None,
                cell_color=None,
                font_size=10):
    """Build a proper PPTX table."""
    hdr_color      = hdr_color      or RED_ROW
    hdr_text_color = hdr_text_color or WHITE
    stripe_color   = stripe_color   or STRIPE
    cell_color     = cell_color     or CARD

    num_cols = len(headers)
    num_rows = 1 + len(rows)

    tbl = sl.shapes.add_table(
        num_rows, num_cols,
        Inches(x), Inches(y), Inches(w), Inches(h)
    ).table

    # column widths
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(w * cw / total)

    def set_cell(cell, text, bold=False,
                 bg=None, fg=None, sz=font_size, align=PP_ALIGN.LEFT):
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg or cell_color
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        if p.runs:
            r = p.runs[0]
        else:
            r = p.add_run()
            r.text = text
        r.font.bold  = bold
        r.font.size  = Pt(sz)
        r.font.color.rgb = fg or WHITE
        r.font.name  = "Calibri"
        cell.margin_left  = Inches(0.08)
        cell.margin_right = Inches(0.05)
        cell.margin_top    = Inches(0.04)
        cell.margin_bottom = Inches(0.04)

    for ci, h in enumerate(headers):
        set_cell(tbl.cell(0, ci), h, bold=True,
                 bg=hdr_color, fg=hdr_text_color, sz=font_size + 0.5,
                 align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        bg = stripe_color if ri % 2 == 1 else cell_color
        for ci, val in enumerate(row):
            set_cell(tbl.cell(ri + 1, ci), str(val), bg=bg)

    return tbl


def bullet_list(sl, items, x, y, w, h,
                size=12, color=None, bullet_color=None, gap=0.27):
    color        = color        or WHITE
    bullet_color = bullet_color or RED
    for i, item in enumerate(items):
        cy = y + i * gap
        txt(sl, "▸", x, cy, 0.22, 0.3, size=size, color=bullet_color, bold=True)
        txt(sl, item, x + 0.22, cy, w - 0.22, 0.32, size=size, color=color)


def mini_card(sl, icon, title, body, x, y, w=2.95, h=1.35,
              accent_color=None, title_color=None):
    accent_color = accent_color or RED
    title_color  = title_color  or RED_LT
    card_box(sl, x, y, w, h, accent_color=accent_color)
    txt(sl, icon,  x + 0.18, y + 0.10, 0.5,       0.38, size=22, color=accent_color)
    txt(sl, title, x + 0.18, y + 0.48, w - 0.28,  0.28, size=11, bold=True, color=title_color)
    txt(sl, body,  x + 0.18, y + 0.76, w - 0.28,  0.55, size=9.5, color=MUTED)


def stat_pill(sl, value, label, x, y, w=2.6, h=0.7,
              val_color=None, pill_color=None):
    val_color  = val_color  or RED_LT
    pill_color = pill_color or CARD2
    card_box(sl, x, y, w, h, accent=False)
    txt(sl, value, x + 0.15, y + 0.04, 1.1, 0.38, size=22, bold=True, color=val_color)
    txt(sl, label, x + 0.15, y + 0.40, w - 0.25, 0.25, size=9, color=MUTED)


def kpi_card(sl, number, label, x, y, w=2.0, h=1.0,
             num_color=None, border_color=None):
    num_color    = num_color    or RED_LT
    border_color = border_color or RED
    card_box(sl, x, y, w, h, accent_color=border_color)
    txt(sl, number, x + 0.12, y + 0.06, w - 0.2, 0.48,
        size=28, bold=True, color=num_color)
    txt(sl, label,  x + 0.12, y + 0.56, w - 0.2, 0.38,
        size=9, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()

# ambient glow ellipse top-left
glow1 = sl.shapes.add_shape(9, Inches(-1.5), Inches(-1.2), Inches(6), Inches(4))
glow1.fill.solid(); glow1.fill.fore_color.rgb = RGBColor(0x22, 0x04, 0x04)
glow1.line.fill.background(); glow1.shadow.inherit = False

# ambient glow ellipse bottom-right
glow2 = sl.shapes.add_shape(9, Inches(9.5), Inches(4.5), Inches(5.5), Inches(4))
glow2.fill.solid(); glow2.fill.fore_color.rgb = RGBColor(0x18, 0x03, 0x03)
glow2.line.fill.background(); glow2.shadow.inherit = False

# shield icon (unicode, large)
txt(sl, "🛡", 5.9, 0.55, 1.5, 1.3, size=54, align=PP_ALIGN.CENTER, color=RED)

# main title
txt(sl, "RiskWaiver360",
    1.5, 1.75, 10.3, 1.2,
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# red underline
divider(sl, 3.2, 2.9, 6.9, color=RED, thick=Pt(2.5))

# subtitle
txt(sl, "GRC Process Exception & Policy Waiver Management Platform",
    1.2, 3.0, 10.9, 0.75,
    size=18, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# track badge
badge = box(sl, 4.4, 3.95, 4.55, 0.45, fill=RED_DIM, border=RED, border_w=Pt(0.75))
txt(sl, "Track: Policy Governance & Risk Management",
    4.4, 3.96, 4.55, 0.44,
    size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# team info placeholder
txt(sl, "Team Name  ·  Institution  ·  2025",
    3.5, 4.65, 6.3, 0.35,
    size=12, color=MUTED2, align=PP_ALIGN.CENTER)

# bottom red line
divider(sl, 0.35, 7.05, 12.65, color=RED, thick=Pt(0.5))

add_notes(sl,
"Welcome. We are presenting RiskWaiver360 — a GRC Process Exception and Policy Waiver Management Platform built for the Policy Governance and Risk Management track.\n\nWe built this to solve a problem every enterprise security team faces but almost none has properly solved: how do you manage security policy exceptions in a way that is centralized, risk-scored, and audit-ready?\n\nLet me start with the problem.")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, 2, "The Problem No One Is Tracking",
               "Policy exceptions scattered, untracked, and audit-invisible")

# context paragraph (left)
txt(sl, "In most enterprises, security policy exceptions live in emails, Excel sheets, and Slack messages — with no central registry, no expiry tracking, and no audit trail.",
    0.35, 1.1, 5.0, 0.95, size=12.5, color=MUTED)

# 4 problem cards (right column, 2×2 grid)
problems = [
    ("📧", "Scattered Records",
     "Exceptions in emails, Excel, Slack. No single source of truth."),
    ("⏰", "Temporary → Permanent",
     "'2-week' waivers stay active for 2 years. No one watches for expiry."),
    ("🔓", "Expired, Never Revoked",
     "Firewall rules, admin access — still live long after the need ends."),
    ("🔍", "Audit Failures",
     "Auditors ask 'who approved what, when?' — teams cannot answer."),
]
positions = [(5.55, 1.05), (9.25, 1.05), (5.55, 3.05), (9.25, 3.05)]
for (icon, title, body), (px, py) in zip(problems, positions):
    mini_card(sl, icon, title, body, px, py, w=3.55, h=1.82)

# screenshot placeholder box
card_box(sl, 0.35, 2.15, 4.85, 2.75, accent_color=MUTED2)
txt(sl, "[ Screenshot: exception tracking chaos\nmockup — inbox / spreadsheet ]",
    0.5, 3.1, 4.55, 0.9, size=10, color=MUTED2, align=PP_ALIGN.CENTER, italic=True)

add_notes(sl,
"In most enterprises, security policy exceptions are not managed — they are tolerated.\n\nWe identified four core failure modes:\n\n1. Scattered records — exceptions live in emails and spreadsheets, no central registry.\n2. Temporary becomes permanent — a 2-week waiver stays active for 2 years.\n3. Expired exceptions never revoked — firewall rules, admin access still live.\n4. Audit failures — when auditors ask for evidence, teams cannot produce it.")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHY IT MATTERS
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, 3, "Six Hidden Risks That Compound Into Breaches",
               "Each risk pattern is detectable — if you have a system that looks for it")

risks = [
    ("👤", "Orphaned Access",
     "Owner left. Exception is still active.\nNo one will renew or revoke it.",
     RED),
    ("📋", "Overlapping Waivers",
     "Two teams, same asset, same policy,\nboth active. Which controls apply?",
     RED_LT),
    ("⚔️",  "Conflicting Approvals",
     "One approver approved. One rejected.\nNobody detected the conflict.",
     AMBER),
    ("📅", "Long-Running Temp Access",
     "Granted for a 2022 project.\nProject ended. Access stayed.",
     AMBER),
    ("📈", "Risk Accumulation",
     "10 Low exceptions on one system\n= Critical combined exposure.",
     RED),
    ("❌", "Compliance / Audit Fail",
     "NIST, GDPR, CIS require documented,\ntime-limited, revocable exceptions.",
     RED_LT),
]
cols = [(0.35, 1.1), (4.55, 1.1), (8.72, 1.1),
        (0.35, 3.45), (4.55, 3.45), (8.72, 3.45)]
for (icon, title, body, ac), (px, py) in zip(risks, cols):
    mini_card(sl, icon, title, body, px, py, w=3.95, h=2.1, accent_color=ac)

# stat row
divider(sl, 0.35, 6.55, 12.65, color=BORDER)
stats = [
    ("68%", "of enterprises cannot produce a complete exception list on demand"),
    ("43%", "of policy exceptions have no recorded expiry date"),
    ("1 in 3", "audit findings relate to access not revoked when it should have been"),
]
for i, (v, l) in enumerate(stats):
    sx = 0.5 + i * 4.35
    txt(sl, v, sx, 6.65, 1.5, 0.45, size=20, bold=True, color=RED_LT)
    txt(sl, l, sx + 1.45, 6.68, 2.7, 0.4, size=9, color=MUTED2, italic=True)

add_notes(sl,
"Let me make this concrete with six specific risk patterns we have modeled in our system.\n\nOrphaned access — an employee who owned an exception has left. The exception is still active. Nobody will revoke it.\n\nOverlapping waivers — two teams submitted an exception for the same asset and policy. Both active. Which compensating controls apply?\n\nConflicting approvals — one approver approved a firewall bypass. A different approver rejected an identical one. Inconsistency that fails governance audits.\n\nLong-running temporary access — vendor granted network access for a project. Project completed. Access never revoked.\n\nRisk accumulation — ten individually Low exceptions on one system represent Critical combined exposure.\n\nCompliance failure — NIST AC-2, GDPR Article 25, and CIS Controls all require documented, justified, time-limited exceptions with evidence of review.")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PROPOSED SOLUTION
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, 4, "RiskWaiver360 — One Platform. Complete Control.",
               "Centralized · Risk-Scored · Lifecycle-Tracked · Audit-Ready")

# headline card
card_box(sl, 0.35, 1.1, 12.65, 1.1)
txt(sl, "RiskWaiver360 centralizes every policy exception into a single risk-scored, lifecycle-tracked, alert-monitored, conflict-detected, and audit-ready governance platform.",
    0.6, 1.22, 12.15, 0.85, size=13.5, color=WHITE, italic=True)

# 5 capability pillars
pillars = [
    ("🗂️",  "Exception\nRegistry",
     "Centralized register.\nSearch, filter, export.",   RED),
    ("🧮", "Risk\nScoring",
     "0–100 rule-based score.\nExplainable per exception.", RED_LT),
    ("🔄", "Lifecycle\nTracking",
     "14-state machine.\nImmutable audit log.",           AMBER),
    ("🧠", "GRC\nIntelligence",
     "Overlap + conflict\ndetection.",                    GREEN),
    ("📊", "Audit\nReport",
     "One-click audit-ready\nevidence report.",           CYAN),
]
px = 0.35
for icon, title, body, ac in pillars:
    card_box(sl, px, 2.45, 2.4, 2.5, accent_color=ac)
    txt(sl, icon,  px + 0.2,  2.55, 0.6, 0.55, size=26, color=ac)
    txt(sl, title, px + 0.2,  3.10, 2.0, 0.55, size=12, bold=True, color=ac)
    txt(sl, body,  px + 0.2,  3.66, 2.0, 0.65, size=9.5, color=MUTED)
    px += 2.52

# screenshot placeholder
card_box(sl, 0.35, 5.15, 12.65, 1.9, accent_color=MUTED2)
txt(sl, "[ Screenshot: Dashboard — KPI cards, charts, risk heatmap ]",
    4.5, 5.85, 5.0, 0.5, size=11, color=MUTED2, align=PP_ALIGN.CENTER, italic=True)

add_notes(sl,
"RiskWaiver360 addresses all the problems we described through five integrated capabilities.\n\nCentralized registry — every exception from submission to revocation in one place, searchable, filterable, exportable, with risk scores and CIA impact visible at a glance.\n\nRisk scoring — automatic 0-to-100 score from seven factors. Explainable — every point traces back to a specific factor.\n\nLifecycle tracking — 14-state machine from draft to closed. Every state transition is logged with actor, timestamp, and reason.\n\nGRC intelligence — cross-portfolio detection of overlaps, conflicting approvals, duplicates, and risk accumulation hotspots.\n\nAudit-ready evidence — one-click report with executive summary, compliance roll-ups, conflict findings, and complete evidence timeline.")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, 5, "Architecture — Full-Stack Working Prototype",
               "React + Vite  ·  Node.js + Express  ·  JSON Demo Storage")

layers = [
    ("USER ROLES",
     "Requester  ·  Security Reviewer  ·  Approver  ·  Auditor / Admin",
     PURPLE, PURPLE_LT),
    ("React + Vite Frontend  (:5173)",
     "Dashboard · Registry · Add Exception · Review Queue · Risk Scoring · Alerts · GRC Intelligence · Audit Report",
     CYAN, CYAN_LT),
    ("Node.js + Express Backend  (:4000)",
     "Risk Scoring Engine · Alert Engine · Conflict Detection · Recommendation Engine · Exception Service · Audit Service",
     RED, RED_LT),
    ("JSON Demo Storage  (backend/data/)",
     "exceptions · approvals · history · alerts · users · policies",
     AMBER, AMBER_LT),
]

y = 1.1
for (title, body, ac, tc) in layers:
    card_box(sl, 0.35, y, 12.65, 1.1, accent_color=ac)
    txt(sl, title, 0.6, y + 0.06, 4.5, 0.38, size=13, bold=True, color=tc)
    txt(sl, body,  0.6, y + 0.46, 12.1, 0.52, size=10.5, color=MUTED)
    # arrow
    if y < 4.5:
        txt(sl, "↓", 6.3, y + 1.1, 0.8, 0.42, size=18, bold=True,
            color=RED, align=PP_ALIGN.CENTER)
    y += 1.52

# tech badge row
badges = ["React 18", "Vite 5", "Node.js 24", "Express 4",
          "Recharts", "lucide-react", "JSON Storage"]
bx = 0.35
for b in badges:
    bw = len(b) * 0.085 + 0.55
    bg_b = box(sl, bx, 6.85, bw, 0.38, fill=RGBColor(0x14, 0x0A, 0x0A), border=RED, border_w=Pt(0.5))
    txt(sl, b, bx + 0.1, 6.88, bw - 0.12, 0.3, size=9.5, bold=True, color=RED_LT)
    bx += bw + 0.18

add_notes(sl,
"The system is a standard full-stack web application. React frontend on Vite at port 5173, Node.js Express backend at port 4000.\n\nEvery screen in the frontend fetches live data from the backend via REST APIs — no hardcoded mock data anywhere in the frontend.\n\nThe backend has six service modules: risk scoring, alert engine, conflict detection, recommendation engine, exception lifecycle service, and audit service.\n\nFor storage we use JSON files — one file per data collection. This means the system runs on any machine with Node.js installed, no database setup required. The architecture is designed so replacing the JSON store with PostgreSQL is a configuration change, not a rewrite.\n\nFor demo authentication, roles are stored in localStorage. In production, this is replaced by SSO and server-side RBAC.")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ROLE-BASED WORKFLOW
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, 6, "Formal Governance Workflow — Four Roles, One System",
               "Every action logged · Immutable timeline · Full role separation")

roles = [
    ("Requester",         PURPLE, PURPLE_LT,
     ["Submit exception via\nGRC Intake Form",
      "Risk auto-scored\non submission",
      "Track own requests\nand alerts"]),
    ("Security Reviewer", RED,    RED_LT,
     ["Triage submitted\nrequests",
      "Add comments,\nforward to approver",
      "See risk score\n& CIA impact"]),
    ("Approver",          AMBER,  AMBER_LT,
     ["Approve / Reject\nRenew / Revoke",
      "Escalate to\nsenior approver",
      "All actions are\nlogged with reason"]),
    ("Auditor / Admin",   GREEN,  GREEN_LT,
     ["Full read-only\nvisibility",
      "Audit Report:\nPrint / Export",
      "GRC Intelligence\n& compliance roll-up"]),
]

rx = 0.35
for role, ac, tc, actions in roles:
    card_box(sl, rx, 1.1, 3.08, 3.9, accent_color=ac)
    txt(sl, role, rx + 0.18, 1.22, 2.7, 0.38, size=13, bold=True, color=tc)
    divider(sl, rx + 0.18, 1.62, 2.7, color=ac, thick=Pt(0.5))
    for ai, action in enumerate(actions):
        ay = 1.75 + ai * 0.95
        txt(sl, "▸", rx + 0.18, ay, 0.22, 0.28, size=10, color=ac, bold=True)
        txt(sl, action, rx + 0.40, ay, 2.42, 0.55, size=10, color=MUTED)
    rx += 3.22

# lifecycle states
divider(sl, 0.35, 5.2, 12.65, color=BORDER)
txt(sl, "Exception Lifecycle States", 0.35, 5.3, 4.0, 0.35,
    size=11, bold=True, color=MUTED)

states = [
    ("DRAFT", MUTED2), ("SUBMITTED", PURPLE_LT), ("UNDER REVIEW", CYAN_LT),
    ("PENDING APPROVAL", AMBER_LT), ("APPROVED", GREEN_LT),
    ("ACTIVE", GREEN), ("EXPIRING SOON", AMBER), ("OVERDUE", RED_LT),
    ("REVOKED", RED), ("CLOSED", MUTED2),
]
sx = 0.35
for state, sc in states:
    sw = len(state) * 0.09 + 0.45
    bg_s = box(sl, sx, 5.72, sw, 0.4, fill=CARD2, border=sc, border_w=Pt(0.6))
    txt(sl, state, sx + 0.1, 5.75, sw - 0.12, 0.32, size=8.5, bold=True, color=sc)
    sx += sw + 0.12
    if sx > 12.5: break
    txt(sl, "→", sx - 0.1, 5.78, 0.2, 0.28, size=9, color=MUTED2)

add_notes(sl,
"The platform models the actual workflow a GRC team would follow.\n\nRequester fills out the GRC Intake Form — selecting the policy, affected asset, exception type, business justification, dates, and compensating controls. Risk is scored automatically on submission.\n\nSecurity Reviewer triages submitted requests — can add comments, ask questions, and forward to the approver.\n\nApprover makes the formal decision: approve, reject, request info, escalate, renew, or revoke. Every action is recorded with identity, timestamp, and note.\n\nAuditor has full read-only access — complete registry, evidence files, and the audit report.\n\nThe state machine at the bottom shows the full lifecycle. Nothing reaches Active without a formal approval record, and every state change is immutable.")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CORE FEATURES
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, 7, "9 Core Features — Production-Ready Scope",
               "Every feature shown is live and working in the prototype")

features = [
    ("📋", "Exception\nRisk Register",
     "Searchable, filterable, CSV\nimport/export. Risk score badge\non every row.",          RED),
    ("📝", "GRC Intake\nForm",
     "Structured multi-step form.\nRisk auto-scored on backend\nat submission time.",        RED_LT),
    ("✅", "Security\nReview Queue",
     "Prioritized worklist. Approve,\nreject, comment. Every action\nis logged.",            AMBER),
    ("🗂️",  "Exception\nEvidence File",
     "Full lifecycle timeline per\nexception. Risk breakdown,\nCIA impact, related alerts.", GREEN),
    ("🧮", "Risk Scoring\nEngine",
     "Formula page: live weights,\nfactor breakdown, ranked list\nof all exceptions.",       CYAN),
    ("🔔", "Continuous\nMonitoring",
     "7 alert types, re-evaluated\non every request. Email\nreminder simulation.",          PURPLE),
    ("🧠", "GRC\nIntelligence",
     "Overlap, conflict, duplicate,\nhotspot detection across\nthe full portfolio.",         RED),
    ("📊", "Audit\nReport",
     "Exec summary, compliance\nroll-ups, conflict findings.\nPrint / export.",             AMBER),
    ("⚙️",  "Settings &\nDemo Config",
     "Role switcher, scoring\nreference table, system\nstatus, user table.",                GREEN),
]

positions = [
    (0.35, 1.12), (4.55, 1.12), (8.72, 1.12),
    (0.35, 3.15), (4.55, 3.15), (8.72, 3.15),
    (0.35, 5.18), (4.55, 5.18), (8.72, 5.18),
]
for (icon, title, body, ac), (px, py) in zip(features, positions):
    mini_card(sl, icon, title, body, px, py, w=3.95, h=1.82, accent_color=ac)

add_notes(sl,
"Let me briefly walk through the nine core features.\n\nException Risk Register — central table with risk score, status badge, CIA impact badge, conflict flag, full-text search, multi-level filtering, CSV import from existing spreadsheets, and CSV export.\n\nGRC Intake Form — structured multi-step form. Risk scored automatically on the backend at submission.\n\nReview Queue — prioritized worklist for reviewers and approvers. Every action recorded in the exception's audit timeline.\n\nEvidence File — full record for one exception: metadata, risk score with point-by-point breakdown, CIA impact, lifecycle history timeline, related alerts and conflict findings.\n\nRisk Scoring page — interactive formula, full weight reference table, ranked list of all exceptions with individual factor breakdowns.\n\nAlerts — seven alert types, continuously re-evaluated. Simulated email reminder preview.\n\nGRC Intelligence — overlap, conflict, duplicate, and hotspot detection.\n\nAudit Report — executive summary, compliance roll-ups, conflict findings, print/export.\n\nSettings — demo configuration, scoring reference, system status.")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — RISK SCORING ENGINE
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, 8, "Rule-Based, Explainable Risk Scoring — 0 to 100",
               "Every point is accountable · No black box · Auditable by design")

# formula card
card_box(sl, 0.35, 1.1, 6.4, 5.35)

formula_lines = [
    "Risk Score  =",
    "",
    "  Exception Type Weight",
    "  + Asset Criticality Weight",
    "  + Duration Penalty",
    "  + Expiry Status Penalty",
    "  + Review Status Penalty",
    "  + Owner Status Penalty",
    "  − Compensating Control Bonus",
    "",
    "  ─────────────────────────────",
    "  Result capped at 0 – 100",
]
txt_block(sl, formula_lines, 0.55, 1.22, 6.0, 5.0,
          size=12.5, color=WHITE, font="Courier New")

# factor weights table
table_shape(sl,
    ["Factor", "Condition", "Points"],
    [
        ["Exception Type",   "Admin Access / Encryption Off",   "+35"],
        ["Exception Type",   "Firewall Bypass",                 "+30"],
        ["Asset Criticality","Critical",                        "+25"],
        ["Asset Criticality","High",                            "+18"],
        ["Duration",         "> 180 days",                      "+20"],
        ["Expiry Status",    "Overdue",                         "+15"],
        ["Expiry Status",    "Expiring ≤ 7 days",               "+8 "],
        ["Review Status",    "Review overdue",                  "+12"],
        ["Owner Status",     "Orphaned owner",                  "+12"],
        ["Comp. Control",    "Control in place",                "−15"],
    ],
    x=6.9, y=1.1, w=6.1, h=4.15,
    col_widths=[2.2, 2.5, 1.0],
    font_size=9.5
)

# risk band bar
band_labels = [("LOW\n0–30", GREEN, 0.35),
               ("MEDIUM\n31–60", AMBER, 3.18),
               ("HIGH\n61–80", RED_LT, 2.4),
               ("CRITICAL\n81–100", RED, 2.45)]
bx = 6.9
for label, c, bw in band_labels:
    box(sl, bx, 5.42, bw, 0.62, fill=c)
    txt(sl, label, bx + 0.05, 5.42, bw - 0.08, 0.62,
        size=8.5, bold=True, color=BG, align=PP_ALIGN.CENTER)
    bx += bw + 0.02

# tag line
txt(sl, "Rule-based  ·  Explainable  ·  Auditable  ·  Not black-box AI",
    6.9, 6.2, 6.1, 0.4, size=11, bold=True, color=RED_LT, align=PP_ALIGN.CENTER)

add_notes(sl,
"The risk scoring model is the analytical core of the platform. Let me explain how it works.\n\nEach exception receives a score between 0 and 100, calculated from seven weighted factors.\n\nException Type — the inherent risk of the exception type. Encryption Disabled and Admin Access carry 35 points each. These weights reflect the potential impact if exploited.\n\nAsset Criticality — how sensitive is the asset? A Critical asset adds 25 points.\n\nDuration — longer exceptions are less likely to be actively monitored. Over 180 days adds 20 points.\n\nExpiry Status — past expiry date adds 15 points. Expiring within 7 days adds 8.\n\nReview Status — overdue review adds 12 points.\n\nOwner Status — orphaned owner adds 12 points.\n\nCompensating Control — if a mitigating control is in place, we deduct 15 points.\n\nThe result maps to four bands: Low 0-30, Medium 31-60, High 61-80, Critical 81-100.\n\nEvery point is traceable. This is deliberate — rule-based, explainable, auditable. Not black-box AI.")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ALERTS & GRC INTELLIGENCE
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, 9, "Continuous Monitoring + Cross-Exception Intelligence",
               "Alert engine evaluates every exception on every request · GRC Intelligence detects portfolio-level patterns")

# left: alert table
txt(sl, "Alert Engine — 7 Types", 0.35, 1.1, 5.8, 0.38,
    size=13, bold=True, color=RED_LT)
table_shape(sl,
    ["Severity", "Alert Type"],
    [
        ["🔴  Critical",  "Exception past expiry date — not yet revoked"],
        ["🔴  Critical",  "Active exception with Critical risk score (≥81)"],
        ["🟠  High",      "Exception expiring within 7 days"],
        ["🟠  High",      "Owner is inactive or has left the organization"],
        ["🟡  Medium",    "No compensating control documented"],
        ["🟡  Medium",    "Exception not reviewed within required period"],
        ["🟢  Low",       "Justification appears vague or generic"],
    ],
    x=0.35, y=1.55, w=6.3, h=3.7,
    col_widths=[1.5, 4.8],
    font_size=10
)

# right: GRC intelligence panels
txt(sl, "GRC Intelligence — 4 Detection Categories", 6.9, 1.1, 6.1, 0.38,
    size=13, bold=True, color=RED_LT)

intel_items = [
    ("🔁", "Overlapping Exceptions",
     "Same asset + policy, both active. Which\ncompensating controls apply?",       RED),
    ("⚔️",  "Conflicting Approvals",
     "One approved, one rejected for the same\nscope. Governance inconsistency.",   AMBER),
    ("📄", "Duplicate Waivers",
     "Near-identical requests from different\nteams. Communication breakdown.",      PURPLE),
    ("📈", "Risk Accumulation Hotspots",
     "10 Low exceptions on one system =\nCritical combined exposure.",              GREEN),
]
iy = 1.55
for icon, title, body, ac in intel_items:
    card_box(sl, 6.9, iy, 6.1, 0.82, accent_color=ac)
    txt(sl, icon,  7.08, iy + 0.08, 0.38, 0.35, size=16, color=ac)
    txt(sl, title, 7.5,  iy + 0.06, 2.5, 0.28, size=11, bold=True, color=ac)
    txt(sl, body,  7.5,  iy + 0.36, 5.3, 0.42, size=9.5, color=MUTED)
    iy += 0.93

# screenshot placeholder
card_box(sl, 0.35, 5.55, 6.3, 1.6, accent_color=MUTED2)
txt(sl, "[ Screenshot: Alerts page — severity badges, email sim card ]",
    1.8, 6.1, 3.7, 0.45, size=10, color=MUTED2, align=PP_ALIGN.CENTER, italic=True)

card_box(sl, 6.9, 5.55, 6.1, 1.6, accent_color=MUTED2)
txt(sl, "[ Screenshot: GRC Intelligence — overlap findings ]",
    8.1, 6.1, 3.7, 0.45, size=10, color=MUTED2, align=PP_ALIGN.CENTER, italic=True)

add_notes(sl,
"The platform has two complementary monitoring layers.\n\nThe Alert Engine evaluates every active exception against seven conditions on every API request — live and always current, no background scheduler needed in the prototype.\n\nCritical alerts fire for overdue exceptions and active exceptions with a risk score above 81. High alerts cover orphaned owners and exceptions expiring within 7 days. Medium covers missing compensating controls and review-overdue conditions. Low flags vague justifications.\n\nThe GRC Intelligence layer is what makes RiskWaiver360 more than a tracking system.\n\nOverlap detection finds cases where two exceptions cover the same asset and policy simultaneously.\n\nConflict detection finds cases where the same scope received both an approval and a rejection from different approvers.\n\nDuplicate detection finds near-identical requests from different requesters.\n\nHotspot analysis aggregates exception risk by asset, business unit, owner, and policy — revealing Critical combined exposure.")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — COMPLIANCE & CIA ALIGNMENT
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, 10, "Built for Compliance — NIST · GDPR · CIS",
               "Compliance alignment designed in from the start — not added as a label")

# compliance table left
txt(sl, "Standards Alignment", 0.35, 1.1, 5.8, 0.35,
    size=13, bold=True, color=RED_LT)
table_shape(sl,
    ["Standard", "Requirement", "How RiskWaiver360 Supports It"],
    [
        ["NIST AC-2",    "Admin access review & revocation",
         "High type-weight + review-overdue alerts + full audit trail"],
        ["NIST PL-4",    "Policy deviations documented",
         "Mandatory justification + approver record + immutable history"],
        ["GDPR Art. 25", "Data protection by design",
         "Encryption-Off scores Critical + CIA: Confidentiality flag"],
        ["CIS 1.1",      "Asset inventory & risk context",
         "Every exception linked to asset · hotspot roll-up by asset"],
    ],
    x=0.35, y=1.52, w=6.3, h=3.2,
    col_widths=[1.2, 2.0, 3.1],
    font_size=9.5
)

# CIA triad right
txt(sl, "CIA Triad Mapping", 6.9, 1.1, 6.1, 0.35,
    size=13, bold=True, color=RED_LT)

cia = [
    ("🔒", "CONFIDENTIALITY",
     "Encryption Disabled · Data Access Extension",  CYAN,   CYAN_LT),
    ("✏️",  "INTEGRITY",
     "Admin Access · Privileged Access Extension\nConflicting Approvals", AMBER, AMBER_LT),
    ("🌐", "AVAILABILITY",
     "Firewall Bypass · Network Exposure",           GREEN,  GREEN_LT),
    ("⚡", "MULTIPLE",
     "Password Policy · Audit Logging Disabled",     PURPLE, PURPLE_LT),
]
cy = 1.52
for icon, label, detail, ac, tc in cia:
    card_box(sl, 6.9, cy, 6.1, 0.95, accent_color=ac)
    txt(sl, icon,   7.08, cy + 0.1,  0.42, 0.4,  size=20, color=ac)
    txt(sl, label,  7.55, cy + 0.08, 2.5,  0.32, size=11, bold=True, color=tc)
    txt(sl, detail, 7.55, cy + 0.44, 5.3,  0.45, size=9.5, color=MUTED)
    cy += 1.07

# key message
card_box(sl, 0.35, 5.55, 12.65, 0.92, accent_color=CYAN)
txt(sl, "CIA mapping translates technical exceptions into business risk language that CISOs, auditors, and board members understand.",
    0.6, 5.67, 12.15, 0.65, size=12, color=WHITE, italic=True)

add_notes(sl,
"Compliance alignment was built into the design from the start, not added as a label at the end.\n\nNIST AC-2 requires privileged access to be controlled and periodically reviewed. Admin Access and Privileged Access exception types carry the highest type weights. Review-overdue alerts fire for these. Every approval and revocation is in the audit trail with actor and timestamp.\n\nNIST PL-4 requires policy deviations to be documented and justified. Every exception requires a business justification stored verbatim in the evidence file. The Vague Justification alert flags thin justifications.\n\nGDPR Article 25 requires data protection by design. Encryption-Disabled and Data Access exceptions automatically score Critical and are tagged CIA: Confidentiality.\n\nCIS 1.1 requires asset inventory with risk context. Every exception is linked to a named asset. Risk is rolled up by asset into hotspot analysis.\n\nThe CIA Triad mapping translates each exception type into its impact dimension — making compliance context immediately visible to CISOs and auditors.")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — DEMO FLOW
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, 11, "End-to-End Demo — 13 Steps, 4 Roles",
               "Every screen calls live backend APIs · No hardcoded mock data")

steps = [
    ("①",  "Home Page",            "Landing page — value proposition & role workspaces",  MUTED),
    ("②",  "Login as Requester",   "requester@riskwaiver360.demo / Requester@123",         PURPLE_LT),
    ("③",  "Submit Exception",     "GRC Intake Form — risk auto-scored on submission",     PURPLE_LT),
    ("④",  "Login as Reviewer",    "reviewer@riskwaiver360.demo / Reviewer@123",           CYAN_LT),
    ("⑤",  "Review Request",       "Review Queue — add comment — forward to approver",     CYAN_LT),
    ("⑥",  "Login as Approver",    "approver@riskwaiver360.demo / Approver@123",           AMBER_LT),
    ("⑦",  "Approve / Revoke",     "Exception Details — risk breakdown — lifecycle action",AMBER_LT),
    ("⑧",  "Dashboard",            "KPI cards — 6 charts — risk heatmap — hotspots",       GREEN_LT),
    ("⑨",  "Exception Registry",   "Filter Critical — conflict badge — open evidence file",RED_LT),
    ("⑩",  "Risk Scoring Page",    "Formula — factor breakdown — ranked exception list",   RED_LT),
    ("⑪",  "Alerts",               "7 alert types — severity filter — email sim preview",  AMBER_LT),
    ("⑫",  "GRC Intelligence",     "Overlap / conflict / duplicate / hotspot findings",    GREEN_LT),
    ("⑬",  "Audit Report",         "Exec summary — compliance roll-up — Print / Export",   CYAN_LT),
]

col1 = steps[:7]
col2 = steps[7:]
for col_idx, col in enumerate([col1, col2]):
    cx = 0.35 if col_idx == 0 else 6.65
    for si, (num, title, detail, tc) in enumerate(col):
        sy = 1.1 + si * 0.84
        card_box(sl, cx, sy, 6.05, 0.76, accent=False)
        # number circle
        circ = box(sl, cx + 0.12, sy + 0.17, 0.42, 0.42, fill=CARD2, border=RED, border_w=Pt(0.75))
        txt(sl, num, cx + 0.12, sy + 0.17, 0.42, 0.42,
            size=11, bold=True, color=tc, align=PP_ALIGN.CENTER)
        txt(sl, title,  cx + 0.65, sy + 0.06, 2.1, 0.3, size=11, bold=True, color=tc)
        txt(sl, detail, cx + 0.65, sy + 0.38, 5.25, 0.32, size=9.5, color=MUTED)

add_notes(sl,
"Here is the demo flow — 13 steps across 4 roles.\n\nWe start on the home page to show the value proposition.\n\nThen login as Requester, submit a new exception through the GRC Intake Form. The backend scores the risk automatically.\n\nSwitch to Reviewer, open the Review Queue, add a comment, forward to approver.\n\nSwitch to Approver, open the evidence file, read the risk score breakdown and recommendation, perform a lifecycle action — approve, revoke, or escalate. Watch the timeline update.\n\nShow the Dashboard — portfolio overview with KPI cards, six charts, and the risk heatmap.\n\nOpen the Registry, filter for Critical risk level, look at the conflict badge, open an exception to show the full Evidence File.\n\nVisit Risk Scoring to show the explainable formula and ranked factor breakdown.\n\nOpen Alerts — filter by Critical severity — show the simulated email reminder.\n\nOpen GRC Intelligence to show the overlap, conflict, and hotspot findings.\n\nFinally open the Audit Report — read the executive summary, point out the conflict findings section, click Print/Export.")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — LIMITATIONS & FUTURE SCOPE
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_header(sl, 12, "Prototype Today. Production Path Clear.",
               "Honest about limitations · Architecture designed for each to have a clean production replacement")

# limitations table
txt(sl, "Current Prototype Limitations", 0.35, 1.1, 6.0, 0.38,
    size=13, bold=True, color=RED_LT)
table_shape(sl,
    ["What", "Current State", "Production Path"],
    [
        ["Authentication",  "localStorage demo only",      "SSO / OIDC (Okta, Entra ID)"],
        ["Authorization",   "Frontend role simulation",    "Server-side RBAC per endpoint"],
        ["Storage",         "JSON files — not concurrent", "PostgreSQL / MongoDB"],
        ["Email alerts",    "Rendered preview only",       "Enterprise SMTP / SendGrid"],
        ["PDF export",      "Browser print dialog",        "Server-side PDF with e-sign"],
        ["Audit actor",     "Client-provided role string", "Server-verified session token"],
    ],
    x=0.35, y=1.55, w=6.3, h=3.0,
    col_widths=[1.6, 2.2, 2.5],
    font_size=9.5
)

# roadmap right
txt(sl, "Production Roadmap", 6.9, 1.1, 6.1, 0.38,
    size=13, bold=True, color=RED_LT)

roadmap = [
    ("🔐", "Security",
     "SSO/OIDC · Backend RBAC\nSecure sessions · Server-side authz",   RED),
    ("🗄️",  "Data & Infra",
     "PostgreSQL/MongoDB\nTamper-evident audit logs · Real PDF",         CYAN),
    ("🔗", "Integrations",
     "ServiceNow/Jira · CMDB\nOkta/Entra ID · SIEM",                   AMBER),
    ("🤖", "Future AI  (scope only — not in prototype)",
     "NLP justification analysis · Anomaly detection\nRAG policy assistant · Predictive escalation", MUTED),
]
ry = 1.55
for icon, title, body, ac in roadmap:
    ht = 1.1 if "Future AI" in title else 0.94
    card_box(sl, 6.9, ry, 6.1, ht, accent_color=ac)
    txt(sl, icon,  7.08, ry + 0.1, 0.4,  0.35, size=18, color=ac)
    txt(sl, title, 7.52, ry + 0.06, 4.8, 0.3,  size=10.5, bold=True, color=ac)
    txt(sl, body,  7.52, ry + 0.42, 5.3, 0.55, size=9.5, color=MUTED)
    ry += ht + 0.08

# closing statement banner
card_box(sl, 0.35, 6.42, 12.65, 0.8, accent_color=RED)
txt(sl, "RiskWaiver360 helps enterprises move from scattered exception tracking to centralized,\nrisk-scored, lifecycle-tracked, and audit-ready GRC exception governance.",
    0.6, 6.5, 12.15, 0.65, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_notes(sl,
"We want to be transparent about what this prototype is and is not.\n\nAuthentication is client-side localStorage only. In production, this is replaced by SSO with server-side session management and per-role permission enforcement on every API endpoint.\n\nStorage is JSON files. The production path is PostgreSQL or MongoDB. The service layer is identical — it is a configuration change, not a rewrite.\n\nEmails are a simulated preview. No real email is sent. Production would use an enterprise SMTP gateway.\n\nExport is browser print. Production would use server-side PDF generation with a digital signature option.\n\nOn AI — the current prototype is fully rule-based. We have scoped AI capabilities as future enhancements: NLP-based justification analysis, anomaly detection, a RAG policy assistant, and predictive escalation. These are future scope only. We believe for a governance and audit tool, starting with explainability is the right design choice.\n\nCLOSING: RiskWaiver360 helps enterprises move from scattered exception tracking to centralized, risk-scored, lifecycle-tracked, and audit-ready GRC exception governance.\n\nThank you.")


# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
prs.save(str(OUT))

size_kb = os.path.getsize(OUT) / 1024
print(f"\nPPTX ready : {OUT}")
print(f"Slides     : {len(prs.slides)}")
print(f"Size       : {size_kb:.1f} KB  ({size_kb/1024:.2f} MB)")
